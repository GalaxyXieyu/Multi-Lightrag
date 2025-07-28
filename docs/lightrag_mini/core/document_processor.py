"""
Document processing functionality for LightRAG Mini
"""

import aiofiles
import asyncio
from pathlib import Path
from typing import List, Dict, Any, Optional
from datetime import datetime
import logging

from .storage import DocStatus, DocProcessingStatus
from .operate import chunking_by_token_size, extract_entities
from .utils import TiktokenTokenizer, compute_mdhash_id, get_content_summary

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Document processor for various file formats"""
    
    def __init__(self):
        self.supported_extensions = (
            ".txt", ".md", ".pdf", ".docx", ".pptx", ".xlsx",
            ".rtf", ".odt", ".tex", ".epub", ".html", ".htm",
            ".csv", ".json", ".xml", ".yaml", ".yml", ".log",
            ".conf", ".ini", ".properties", ".sql", ".bat", ".sh",
            ".c", ".cpp", ".py", ".java", ".js", ".ts", ".swift",
            ".go", ".rb", ".php", ".css", ".scss", ".less"
        )
    
    def is_supported_file(self, filename: str) -> bool:
        """Check if file is supported"""
        return any(filename.lower().endswith(ext) for ext in self.supported_extensions)
    
    async def extract_text_from_file(self, file_path: Path) -> str:
        """Extract text content from file based on extension"""
        ext = file_path.suffix.lower()
        
        try:
            if ext in [".txt", ".md", ".html", ".htm", ".tex", ".json", ".xml", 
                      ".yaml", ".yml", ".rtf", ".odt", ".epub", ".csv", ".log",
                      ".conf", ".ini", ".properties", ".sql", ".bat", ".sh",
                      ".c", ".cpp", ".py", ".java", ".js", ".ts", ".swift",
                      ".go", ".rb", ".php", ".css", ".scss", ".less"]:
                return await self._extract_text_file(file_path)
            elif ext == ".pdf":
                return await self._extract_pdf(file_path)
            elif ext == ".docx":
                return await self._extract_docx(file_path)
            elif ext == ".pptx":
                return await self._extract_pptx(file_path)
            elif ext == ".xlsx":
                return await self._extract_xlsx(file_path)
            else:
                raise ValueError(f"Unsupported file type: {ext}")
                
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            raise
    
    async def _extract_text_file(self, file_path: Path) -> str:
        """Extract text from plain text files"""
        async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
            content = await f.read()
            
        if not content or len(content.strip()) == 0:
            raise ValueError(f"Empty content in file: {file_path.name}")
        
        return content
    
    async def _extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF files"""
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            raise ImportError("PyPDF2 is required for PDF processing. Install with: pip install PyPDF2")
        
        reader = PdfReader(str(file_path))
        content = ""
        
        for page in reader.pages:
            content += page.extract_text() + "\n"
        
        if not content.strip():
            raise ValueError(f"No text content found in PDF: {file_path.name}")
        
        return content
    
    async def _extract_docx(self, file_path: Path) -> str:
        """Extract text from DOCX files"""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx is required for DOCX processing. Install with: pip install python-docx")
        
        doc = Document(str(file_path))
        content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        
        if not content.strip():
            raise ValueError(f"No text content found in DOCX: {file_path.name}")
        
        return content
    
    async def _extract_pptx(self, file_path: Path) -> str:
        """Extract text from PPTX files"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx is required for PPTX processing. Install with: pip install python-pptx")
        
        prs = Presentation(str(file_path))
        content = ""
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
        
        if not content.strip():
            raise ValueError(f"No text content found in PPTX: {file_path.name}")
        
        return content
    
    async def _extract_xlsx(self, file_path: Path) -> str:
        """Extract text from XLSX files"""
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise ImportError("openpyxl is required for XLSX processing. Install with: pip install openpyxl")
        
        wb = load_workbook(str(file_path))
        content = ""
        
        for sheet in wb:
            content += f"Sheet: {sheet.title}\n"
            for row in sheet.iter_rows(values_only=True):
                content += "\t".join(str(cell) if cell is not None else "" for cell in row) + "\n"
            content += "\n"
        
        if not content.strip():
            raise ValueError(f"No text content found in XLSX: {file_path.name}")
        
        return content
    
    async def process_file(
        self,
        file_path: Path,
        tokenizer: TiktokenTokenizer,
        llm_func: callable,
        chunk_size: int = 1200,
        chunk_overlap: int = 100,
        cache_storage: Optional[Dict] = None,
        **llm_kwargs
    ) -> Dict[str, Any]:
        """
        Process a file and extract entities/relationships
        
        Args:
            file_path: Path to the file
            tokenizer: Tokenizer for chunking
            llm_func: LLM function for entity extraction
            chunk_size: Maximum tokens per chunk
            chunk_overlap: Overlapping tokens between chunks
            cache_storage: Optional cache storage
            **llm_kwargs: Additional LLM arguments
        
        Returns:
            Dictionary containing processing results
        """
        try:
            # Extract text content
            content = await self.extract_text_from_file(file_path)
            
            # Generate document ID
            doc_id = compute_mdhash_id(content)
            
            # Create document status
            doc_status = DocProcessingStatus(
                content_summary=get_content_summary(content),
                content_length=len(content),
                status=DocStatus.PROCESSING,
                created_at=datetime.now().isoformat(),
                updated_at=datetime.now().isoformat(),
                file_path=str(file_path.name)
            )
            
            # Split into chunks
            chunks = chunking_by_token_size(
                tokenizer=tokenizer,
                content=content,
                max_token_size=chunk_size,
                overlap_token_size=chunk_overlap
            )
            
            # Extract entities and relationships from each chunk
            all_entities = []
            all_relationships = []
            
            for i, chunk in enumerate(chunks):
                chunk_id = f"{doc_id}_chunk_{i}"
                
                try:
                    extraction_result = await extract_entities(
                        llm_func=llm_func,
                        chunk_content=chunk["content"],
                        chunk_id=chunk_id,
                        cache_storage=cache_storage,
                        **llm_kwargs
                    )
                    
                    entities = extraction_result.get("entities", [])
                    relationships = extraction_result.get("relationships", [])
                    
                    # Add chunk info to entities and relationships
                    for entity in entities:
                        entity["chunk_id"] = chunk_id
                        entity["doc_id"] = doc_id
                    
                    for rel in relationships:
                        rel["chunk_id"] = chunk_id
                        rel["doc_id"] = doc_id
                    
                    all_entities.extend(entities)
                    all_relationships.extend(relationships)
                    
                except Exception as e:
                    logger.error(f"Error processing chunk {chunk_id}: {str(e)}")
                    continue
            
            # Update document status
            doc_status.status = DocStatus.PROCESSED
            doc_status.chunks_count = len(chunks)
            doc_status.updated_at = datetime.now().isoformat()
            
            return {
                "doc_id": doc_id,
                "doc_status": doc_status,
                "content": content,
                "chunks": chunks,
                "entities": all_entities,
                "relationships": all_relationships
            }
            
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {str(e)}")
            
            # Create failed status
            doc_status = DocProcessingStatus(
                content_summary="Processing failed",
                content_length=0,
                status=DocStatus.FAILED,
                created_at=datetime.now().isoformat(),
                updated_at=datetime.now().isoformat(),
                error=str(e),
                file_path=str(file_path.name)
            )
            
            return {
                "doc_id": None,
                "doc_status": doc_status,
                "content": "",
                "chunks": [],
                "entities": [],
                "relationships": [],
                "error": str(e)
            }
    
    async def process_text(
        self,
        text: str,
        source: str,
        tokenizer: TiktokenTokenizer,
        llm_func: callable,
        chunk_size: int = 1200,
        chunk_overlap: int = 100,
        cache_storage: Optional[Dict] = None,
        **llm_kwargs
    ) -> Dict[str, Any]:
        """
        Process raw text and extract entities/relationships
        
        Args:
            text: Text content
            source: Source identifier
            tokenizer: Tokenizer for chunking
            llm_func: LLM function for entity extraction
            chunk_size: Maximum tokens per chunk
            chunk_overlap: Overlapping tokens between chunks
            cache_storage: Optional cache storage
            **llm_kwargs: Additional LLM arguments
        
        Returns:
            Dictionary containing processing results
        """
        try:
            # Generate document ID
            doc_id = compute_mdhash_id(text)
            
            # Create document status
            doc_status = DocProcessingStatus(
                content_summary=get_content_summary(text),
                content_length=len(text),
                status=DocStatus.PROCESSING,
                created_at=datetime.now().isoformat(),
                updated_at=datetime.now().isoformat(),
                file_path=source
            )
            
            # Split into chunks
            chunks = chunking_by_token_size(
                tokenizer=tokenizer,
                content=text,
                max_token_size=chunk_size,
                overlap_token_size=chunk_overlap
            )
            
            # Extract entities and relationships
            all_entities = []
            all_relationships = []
            
            for i, chunk in enumerate(chunks):
                chunk_id = f"{doc_id}_chunk_{i}"
                
                try:
                    extraction_result = await extract_entities(
                        llm_func=llm_func,
                        chunk_content=chunk["content"],
                        chunk_id=chunk_id,
                        cache_storage=cache_storage,
                        **llm_kwargs
                    )
                    
                    entities = extraction_result.get("entities", [])
                    relationships = extraction_result.get("relationships", [])
                    
                    # Add metadata
                    for entity in entities:
                        entity["chunk_id"] = chunk_id
                        entity["doc_id"] = doc_id
                    
                    for rel in relationships:
                        rel["chunk_id"] = chunk_id
                        rel["doc_id"] = doc_id
                    
                    all_entities.extend(entities)
                    all_relationships.extend(relationships)
                    
                except Exception as e:
                    logger.error(f"Error processing chunk {chunk_id}: {str(e)}")
                    continue
            
            # Update status
            doc_status.status = DocStatus.PROCESSED
            doc_status.chunks_count = len(chunks)
            doc_status.updated_at = datetime.now().isoformat()
            
            return {
                "doc_id": doc_id,
                "doc_status": doc_status,
                "content": text,
                "chunks": chunks,
                "entities": all_entities,
                "relationships": all_relationships
            }
            
        except Exception as e:
            logger.error(f"Error processing text: {str(e)}")
            
            doc_status = DocProcessingStatus(
                content_summary="Processing failed",
                content_length=0,
                status=DocStatus.FAILED,
                created_at=datetime.now().isoformat(),
                updated_at=datetime.now().isoformat(),
                error=str(e),
                file_path=source
            )
            
            return {
                "doc_id": None,
                "doc_status": doc_status,
                "content": "",
                "chunks": [],
                "entities": [],
                "relationships": [],
                "error": str(e)
            }